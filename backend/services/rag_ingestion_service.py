"""
RAG Ingestion Service — load → chunk → insert into Delta → sync VS index.

Duplicate handling
──────────────────
Every ingest call first computes SHA-256 of the raw file bytes (doc_hash).
Two duplicate scenarios are detected before any INSERT:

  exact_content   — doc_hash already exists in the table (identical bytes,
                    possibly under a different filename)
  same_filename   — source column matches filename but doc_hash differs
                    (updated version of the same document)

When a duplicate is detected the service returns
  { action_required: True, duplicate_type: "...", ... }
without writing anything.  The caller then re-submits with mode=:

  "replace"  — DELETE all rows WHERE source = filename, then INSERT fresh chunks
  "append"   — INSERT new chunks alongside existing ones (different chunk_ids)
  "skip"     — do nothing, return early

Chunk IDs are deterministic: SHA-256(source :: chunk_index :: content[:128])
truncated to 40 hex chars.  Re-ingesting the same document with replace
produces identical chunk_ids, so the VS index converges without ghost vectors.

Delta table schema adds:  doc_hash STRING
"""

from __future__ import annotations

import hashlib
import io
import json
import logging
from datetime import datetime, timezone
from typing import Any, Literal

from backend.integrations.databricks.sql_statements import execute_sql_statement
from backend.integrations.databricks.vector_search import (
    ensure_vs_endpoint,
    ensure_delta_sync_index,
    sync_index,
    SyncNotReady,
)

logger = logging.getLogger(__name__)

# ── Defaults ──────────────────────────────────────────────────────────────────

CATALOG = "workspace"
SCHEMA = "agentic_rag"
DEFAULT_TABLE = "documents"
DEFAULT_VS_ENDPOINT = "rag_endpoint"
DEFAULT_INDEX = "rag_index"
DEFAULT_PRIMARY_KEY = "chunk_id"
DEFAULT_CONTENT_COL = "content"
DEFAULT_EMBEDDING_ENDPOINT = "databricks-bge-large-en"

IngestMode = Literal["check", "replace", "append", "skip"]

# ── DDL ───────────────────────────────────────────────────────────────────────

_CREATE_TABLE_SQL = """
CREATE TABLE IF NOT EXISTS {fqn} (
  chunk_id     STRING  NOT NULL,
  source       STRING,
  source_type  STRING,
  content      STRING,
  page         INT,
  section      STRING,
  author       STRING,
  ingested_at  TIMESTAMP,
  doc_hash     STRING
)
USING DELTA
TBLPROPERTIES ('delta.enableChangeDataFeed' = 'true')
"""

# ADD COLUMN is idempotent-safe via IF NOT EXISTS (Databricks Delta supports this)
_ADD_DOC_HASH_COL_SQL = """
ALTER TABLE {fqn} ADD COLUMN IF NOT EXISTS doc_hash STRING
"""


def ensure_delta_table(
    host: str,
    token: str,
    warehouse_id: str,
    catalog: str,
    schema: str,
    table: str,
) -> dict[str, Any]:
    fqn = f"`{catalog}`.`{schema}`.`{table}`"
    result = execute_sql_statement(host, token, warehouse_id, _CREATE_TABLE_SQL.format(fqn=fqn))
    if not result["ok"]:
        raise RuntimeError(f"Could not create Delta table {fqn}: {result.get('error')}")
    # Ensure doc_hash column exists on older tables that predate this migration
    execute_sql_statement(host, token, warehouse_id, _ADD_DOC_HASH_COL_SQL.format(fqn=fqn))
    logger.info("Delta table %s ensured", fqn)
    return result


# ── Duplicate detection ───────────────────────────────────────────────────────

def _file_hash(file_bytes: bytes) -> str:
    return hashlib.sha256(file_bytes).hexdigest()


def check_duplicate(
    host: str,
    token: str,
    warehouse_id: str,
    fqn: str,
    doc_hash: str,
    filename: str,
) -> dict[str, Any]:
    """
    Returns:
      { duplicate_type: None | "exact_content" | "same_filename",
        existing_source: str | None,
        existing_chunks: int,
        doc_hash: str }

    exact_content  — same bytes already in the table (maybe different name)
    same_filename  — same source name, different content (updated doc)
    None           — new document, safe to ingest
    """
    escaped_hash = doc_hash.replace("'", "\\'")
    escaped_name = filename.replace("'", "\\'")

    # Check 1: same content hash (exact duplicate, any source name)
    hash_sql = (
        f"SELECT source, COUNT(*) AS cnt FROM {fqn} "
        f"WHERE doc_hash = '{escaped_hash}' GROUP BY source LIMIT 5"
    )
    hash_result = execute_sql_statement(host, token, warehouse_id, hash_sql)
    if hash_result["ok"] and hash_result["rows"]:
        existing_sources = [r[0] for r in hash_result["rows"]]
        total_chunks = sum(int(r[1]) for r in hash_result["rows"])
        return {
            "duplicate_type": "exact_content",
            "existing_source": existing_sources[0],
            "all_sources": existing_sources,
            "existing_chunks": total_chunks,
            "doc_hash": doc_hash,
        }

    # Check 2: same filename, different content (updated version)
    name_sql = (
        f"SELECT COUNT(*) AS cnt FROM {fqn} "
        f"WHERE source = '{escaped_name}' AND (doc_hash IS NULL OR doc_hash != '{escaped_hash}')"
    )
    name_result = execute_sql_statement(host, token, warehouse_id, name_sql)
    if name_result["ok"] and name_result["rows"]:
        cnt = int(name_result["rows"][0][0] or 0)
        if cnt > 0:
            return {
                "duplicate_type": "same_filename",
                "existing_source": filename,
                "all_sources": [filename],
                "existing_chunks": cnt,
                "doc_hash": doc_hash,
            }

    return {
        "duplicate_type": None,
        "existing_source": None,
        "all_sources": [],
        "existing_chunks": 0,
        "doc_hash": doc_hash,
    }


def delete_chunks_by_source(
    host: str,
    token: str,
    warehouse_id: str,
    fqn: str,
    filename: str,
) -> int:
    """Delete all chunks for a given source filename. Returns deleted count (best-effort)."""
    escaped = filename.replace("'", "\\'")
    # Count first so we can report it
    count_sql = f"SELECT COUNT(*) FROM {fqn} WHERE source = '{escaped}'"
    count_res = execute_sql_statement(host, token, warehouse_id, count_sql)
    deleted = 0
    if count_res["ok"] and count_res["rows"]:
        deleted = int(count_res["rows"][0][0] or 0)
    delete_sql = f"DELETE FROM {fqn} WHERE source = '{escaped}'"
    del_res = execute_sql_statement(host, token, warehouse_id, delete_sql)
    if not del_res["ok"]:
        raise RuntimeError(f"DELETE failed: {del_res.get('error')}")
    logger.info("Deleted %d existing chunks for source '%s'", deleted, filename)
    return deleted


# ── Document loading ──────────────────────────────────────────────────────────

def _load_pdf(file_bytes: bytes) -> list[dict[str, Any]]:
    try:
        import PyPDF2
        reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
        pages = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            if text.strip():
                pages.append({"text": text, "page": i + 1, "source_type": "pdf"})
        return pages
    except Exception as exc:
        raise RuntimeError(f"PDF load failed: {exc}") from exc


def _load_html(file_bytes: bytes) -> list[dict[str, Any]]:
    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(file_bytes, "html.parser")
        text = soup.get_text(separator="\n", strip=True)
        return [{"text": text, "page": 1, "source_type": "html"}]
    except Exception as exc:
        raise RuntimeError(f"HTML load failed: {exc}") from exc


def _load_json(file_bytes: bytes) -> list[dict[str, Any]]:
    try:
        data = json.loads(file_bytes.decode("utf-8", errors="replace"))
        text = json.dumps(data, indent=2, ensure_ascii=False)
        return [{"text": text, "page": 1, "source_type": "json"}]
    except Exception as exc:
        raise RuntimeError(f"JSON load failed: {exc}") from exc


def _load_text(file_bytes: bytes) -> list[dict[str, Any]]:
    text = file_bytes.decode("utf-8", errors="replace")
    return [{"text": text, "page": 1, "source_type": "text"}]


def _load_csv(file_bytes: bytes) -> list[dict[str, Any]]:
    import csv as _csv
    try:
        text = file_bytes.decode("utf-8", errors="replace")
        reader = _csv.reader(io.StringIO(text))
        rows = list(reader)
        if not rows:
            return [{"text": "", "page": 1, "source_type": "csv"}]
        headers = rows[0]
        lines = []
        for row in rows[1:]:
            pairs = [f"{h}: {v}" for h, v in zip(headers, row) if str(v).strip()]
            if pairs:
                lines.append("  |  ".join(pairs))
        return [{"text": "\n".join(lines), "page": 1, "source_type": "csv"}]
    except Exception as exc:
        raise RuntimeError(f"CSV load failed: {exc}") from exc


def _load_xml(file_bytes: bytes) -> list[dict[str, Any]]:
    import xml.etree.ElementTree as ET
    try:
        root = ET.fromstring(file_bytes.decode("utf-8", errors="replace"))
        texts = [el.text.strip() for el in root.iter() if el.text and el.text.strip()]
        return [{"text": "\n".join(texts), "page": 1, "source_type": "xml"}]
    except Exception as exc:
        raise RuntimeError(f"XML load failed: {exc}") from exc


def _load_excel(file_bytes: bytes) -> list[dict[str, Any]]:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
        pages = []
        for sheet_num, ws in enumerate(wb.worksheets, 1):
            lines = []
            headers: list[str] = []
            for row_idx, row in enumerate(ws.iter_rows(values_only=True)):
                vals = [str(c) if c is not None else "" for c in row]
                if row_idx == 0:
                    headers = vals
                    continue
                pairs = (
                    [f"{h}: {v}" for h, v in zip(headers, vals) if v.strip()]
                    if headers else [v for v in vals if v.strip()]
                )
                if pairs:
                    lines.append("  |  ".join(pairs))
            if lines:
                pages.append({"text": "\n".join(lines), "page": sheet_num, "source_type": "excel"})
        return pages or [{"text": "", "page": 1, "source_type": "excel"}]
    except Exception as exc:
        raise RuntimeError(f"Excel load failed: {exc}") from exc


def _load_word(file_bytes: bytes) -> list[dict[str, Any]]:
    try:
        from docx import Document
        doc = Document(io.BytesIO(file_bytes))
        lines = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    lines.append("  |  ".join(cells))
        return [{"text": "\n".join(lines), "page": 1, "source_type": "word"}]
    except Exception as exc:
        raise RuntimeError(f"Word load failed: {exc}") from exc


def _load_ppt(file_bytes: bytes) -> list[dict[str, Any]]:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        pages = []
        for slide_num, slide in enumerate(prs.slides, 1):
            texts = [
                shape.text.strip()
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]
            if texts:
                pages.append({"text": "\n".join(texts), "page": slide_num, "source_type": "ppt"})
        return pages or [{"text": "", "page": 1, "source_type": "ppt"}]
    except Exception as exc:
        raise RuntimeError(f"PPT load failed: {exc}") from exc


def load_document(file_bytes: bytes, source_type: str) -> list[dict[str, Any]]:
    st = (source_type or "text").lower()
    if st == "pdf":    return _load_pdf(file_bytes)
    if st == "html":   return _load_html(file_bytes)
    if st == "json":   return _load_json(file_bytes)
    if st == "csv":    return _load_csv(file_bytes)
    if st == "xml":    return _load_xml(file_bytes)
    if st == "excel":  return _load_excel(file_bytes)
    if st == "word":   return _load_word(file_bytes)
    if st == "ppt":    return _load_ppt(file_bytes)
    return _load_text(file_bytes)


# ── Chunking ──────────────────────────────────────────────────────────────────

def _chunk_recursive(text: str, chunk_size: int = 1200, chunk_overlap: int = 300) -> list[str]:
    from langchain.text_splitter import RecursiveCharacterTextSplitter
    return RecursiveCharacterTextSplitter(
        chunk_size=chunk_size, chunk_overlap=chunk_overlap, length_function=len
    ).split_text(text)


def _chunk_fixed(text: str, chunk_size: int = 1200, chunk_overlap: int = 300) -> list[str]:
    from langchain.text_splitter import CharacterTextSplitter
    return CharacterTextSplitter(
        chunk_size=chunk_size, chunk_overlap=chunk_overlap,
        separator="\n\n", length_function=len
    ).split_text(text)


def _chunk_semantic(
    text: str,
    breakpoint_type: str = "percentile",
    breakpoint_threshold: int = 95,
    fallback_chunk_size: int = 1200,
    chunk_overlap: int = 150,
) -> list[str]:
    try:
        from langchain_experimental.text_splitter import SemanticChunker
        from langchain_community.embeddings import HuggingFaceEmbeddings
        embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
        return SemanticChunker(
            embeddings,
            breakpoint_threshold_type=breakpoint_type,
            breakpoint_threshold_amount=breakpoint_threshold,
        ).split_text(text)
    except ImportError:
        logger.warning("langchain_experimental not installed — falling back to recursive splitter")
        return _chunk_recursive(text, fallback_chunk_size, chunk_overlap)
    except Exception as exc:
        logger.warning("SemanticChunker failed (%s) — falling back to recursive splitter", exc)
        return _chunk_recursive(text, fallback_chunk_size, chunk_overlap)


def chunk_document(page: dict[str, Any], chunker_cfg: dict[str, Any]) -> list[str]:
    text = page["text"]
    strategy = (chunker_cfg.get("strategy") or "recursive").lower()
    chunk_size = int(chunker_cfg.get("fallback_chunk_size") or 1200)
    overlap = int(chunker_cfg.get("chunk_overlap") or 300)
    if strategy == "semantic":
        return _chunk_semantic(
            text,
            breakpoint_type=chunker_cfg.get("breakpoint_type") or "percentile",
            breakpoint_threshold=int(chunker_cfg.get("breakpoint_threshold") or 95),
            fallback_chunk_size=chunk_size,
            chunk_overlap=overlap,
        )
    if strategy == "fixed":
        return _chunk_fixed(text, chunk_size, overlap)
    return _chunk_recursive(text, chunk_size, overlap)


# ── Chunk ID (deterministic) ──────────────────────────────────────────────────

def _make_chunk_id(source: str, chunk_index: int, content: str) -> str:
    """
    Deterministic 40-char hex ID based on source name, position, and content prefix.
    Same document re-ingested → same IDs → VS index upserts cleanly.
    """
    raw = f"{source}::{chunk_index}::{content[:128]}"
    return hashlib.sha256(raw.encode("utf-8")).hexdigest()[:40]


# ── Delta write ───────────────────────────────────────────────────────────────

def _escape_sql_str(s: str) -> str:
    return s.replace("\\", "\\\\").replace("'", "\\'")


def _build_insert_sql(
    fqn: str,
    rows: list[dict[str, Any]],
    metadata_enrichment: str = "basic",
) -> str:
    now_ts = datetime.now(timezone.utc).strftime("%Y-%m-%d %H:%M:%S")
    parts: list[str] = []
    for row in rows:
        chunk_id    = _escape_sql_str(str(row.get("chunk_id", "")))
        source      = _escape_sql_str(str(row.get("source", "")))
        source_type = _escape_sql_str(str(row.get("source_type", "")))
        content     = _escape_sql_str(str(row.get("content", "")))
        page        = int(row.get("page") or 0)
        section     = _escape_sql_str(str(row.get("section", ""))) if metadata_enrichment == "extended" else ""
        author      = _escape_sql_str(str(row.get("author", "")))  if metadata_enrichment == "extended" else ""
        doc_hash    = _escape_sql_str(str(row.get("doc_hash", "")))
        parts.append(
            f"('{chunk_id}', '{source}', '{source_type}', '{content}', "
            f"{page}, '{section}', '{author}', TIMESTAMP '{now_ts}', '{doc_hash}')"
        )
    values = ",\n  ".join(parts)
    return (
        f"INSERT INTO {fqn} "
        f"(chunk_id, source, source_type, content, page, section, author, ingested_at, doc_hash)\n"
        f"VALUES\n  {values}"
    )


def insert_chunks(
    host: str,
    token: str,
    warehouse_id: str,
    fqn: str,
    chunks: list[dict[str, Any]],
    metadata_enrichment: str = "basic",
    batch_size: int = 50,
) -> dict[str, Any]:
    total = len(chunks)
    inserted = 0
    errors: list[str] = []
    for i in range(0, total, batch_size):
        batch = chunks[i : i + batch_size]
        sql = _build_insert_sql(fqn, batch, metadata_enrichment)
        result = execute_sql_statement(host, token, warehouse_id, sql)
        if result["ok"]:
            inserted += len(batch)
        else:
            err = result.get("error", "unknown")
            logger.error("INSERT batch %d failed: %s", i // batch_size, err)
            errors.append(err)
    return {"total": total, "inserted": inserted, "errors": errors}


# ── Provision ─────────────────────────────────────────────────────────────────

def provision(
    host: str,
    token: str,
    warehouse_id: str,
    ingestion_cfg: dict[str, Any],
) -> dict[str, Any]:
    db = ingestion_cfg.get("databricks", {})
    catalog          = db.get("catalog")          or CATALOG
    schema           = db.get("schema")           or SCHEMA
    table            = db.get("table")            or DEFAULT_TABLE
    vs_endpoint_name = db.get("vs_endpoint")      or DEFAULT_VS_ENDPOINT
    index_name       = db.get("index_name")       or DEFAULT_INDEX
    primary_key      = db.get("primary_key")      or DEFAULT_PRIMARY_KEY
    content_col      = db.get("content_column")   or DEFAULT_CONTENT_COL
    pipeline_type    = db.get("pipeline_type")    or "TRIGGERED"
    embedding_endpoint = (
        ingestion_cfg.get("embeddings", {}).get("databricks_endpoint")
        or DEFAULT_EMBEDDING_ENDPOINT
    )
    fqn_table        = f"`{catalog}`.`{schema}`.`{table}`"
    index_full_name  = f"{catalog}.{schema}.{index_name}"
    table_full_name  = f"{catalog}.{schema}.{table}"

    table_result = ensure_delta_table(host, token, warehouse_id, catalog, schema, table)
    ensure_vs_endpoint(host, token, vs_endpoint_name)
    index_info = ensure_delta_sync_index(
        host, token,
        endpoint_name=vs_endpoint_name,
        index_full_name=index_full_name,
        source_table_full_name=table_full_name,
        primary_key=primary_key,
        content_column=content_col,
        embedding_endpoint=embedding_endpoint,
        pipeline_type=pipeline_type,
    )
    return {
        "ok": True,
        "table_fqn": fqn_table,
        "vs_endpoint": vs_endpoint_name,
        "index": index_full_name,
        "index_state": index_info.get("status", {}).get("detailed_state", "PROVISIONING"),
        "table_created": table_result.get("ok"),
    }


# ── Main pipeline ─────────────────────────────────────────────────────────────

def run_ingestion(
    host: str,
    token: str,
    warehouse_id: str,
    ingestion_cfg: dict[str, Any],
    file_bytes: bytes,
    filename: str = "upload",
    mode: IngestMode = "check",
) -> dict[str, Any]:
    """
    Full ingestion pipeline with duplicate guard.

    mode="check"   → detect duplicates, return action_required without writing
    mode="replace" → delete existing chunks for this filename, then ingest
    mode="append"  → ingest alongside existing chunks (new chunk_ids)
    mode="skip"    → return early without writing
    """
    db               = ingestion_cfg.get("databricks", {})
    catalog          = db.get("catalog")        or CATALOG
    schema           = db.get("schema")         or SCHEMA
    table            = db.get("table")          or DEFAULT_TABLE
    vs_endpoint_name = db.get("vs_endpoint")    or DEFAULT_VS_ENDPOINT
    index_name       = db.get("index_name")     or DEFAULT_INDEX
    pipeline_type    = db.get("pipeline_type")  or "TRIGGERED"
    metadata_enrichment = ingestion_cfg.get("metadata_enrichment") or "basic"
    source_type      = ingestion_cfg.get("source_type") or "pdf"
    chunker_cfg      = ingestion_cfg.get("chunker") or {}

    fqn_table       = f"`{catalog}`.`{schema}`.`{table}`"
    index_full_name = f"{catalog}.{schema}.{index_name}"

    # ── Compute file hash ──────────────────────────────────────────────────
    doc_hash = _file_hash(file_bytes)

    # ── Skip shortcut ──────────────────────────────────────────────────────
    if mode == "skip":
        return {"ok": True, "skipped": True, "filename": filename, "doc_hash": doc_hash}

    # ── Duplicate check (always run, even in replace/append mode for the log) ──
    dup_info = check_duplicate(host, token, warehouse_id, fqn_table, doc_hash, filename)

    if mode == "check" and dup_info["duplicate_type"] is not None:
        return {
            "ok": False,
            "action_required": True,
            "duplicate_type": dup_info["duplicate_type"],
            "existing_source": dup_info["existing_source"],
            "all_sources": dup_info["all_sources"],
            "existing_chunks": dup_info["existing_chunks"],
            "doc_hash": doc_hash,
            "filename": filename,
        }

    # ── Replace: purge old chunks first ───────────────────────────────────
    deleted_count = 0
    if mode == "replace":
        deleted_count = delete_chunks_by_source(host, token, warehouse_id, fqn_table, filename)

    # ── Load ──────────────────────────────────────────────────────────────
    logger.info("Loading '%s' (source_type=%s, mode=%s)", filename, source_type, mode)
    pages = load_document(file_bytes, source_type)
    if not pages:
        return {"ok": False, "error": "No text extracted from document"}

    # ── Chunk ─────────────────────────────────────────────────────────────
    all_chunks: list[dict[str, Any]] = []
    global_index = 0
    for page in pages:
        texts = chunk_document(page, chunker_cfg)
        for t in texts:
            if not t.strip():
                continue
            all_chunks.append({
                "chunk_id":    _make_chunk_id(filename, global_index, t),
                "source":      filename,
                "source_type": source_type,
                "content":     t,
                "page":        page.get("page") or 1,
                "section":     page.get("section") or "",
                "author":      page.get("author") or "",
                "doc_hash":    doc_hash,
            })
            global_index += 1

    if not all_chunks:
        return {"ok": False, "error": "Chunking produced no non-empty chunks"}

    logger.info("Generated %d chunks from %d page(s)", len(all_chunks), len(pages))

    # ── Insert ────────────────────────────────────────────────────────────
    insert_result = insert_chunks(host, token, warehouse_id, fqn_table, all_chunks, metadata_enrichment)
    if insert_result["errors"]:
        logger.warning("Some INSERT batches failed: %s", insert_result["errors"])

    # ── VS sync ───────────────────────────────────────────────────────────
    sync_triggered = False
    sync_error: str | None = None
    sync_pending: bool = False
    if pipeline_type == "TRIGGERED":
        try:
            sync_index(host, token, index_full_name)
            sync_triggered = True
        except SyncNotReady as exc:
            # Index still provisioning — chunks are safely in Delta,
            # VS will pick them up automatically once it reaches ONLINE state.
            sync_pending = True
            sync_error = "Index still initializing — will sync automatically once ONLINE"
            logger.info("VS sync deferred (index not ready): %s", exc)
        except Exception as exc:
            sync_error = str(exc)
            logger.warning("VS sync trigger failed: %s", exc)

    return {
        "ok": insert_result["inserted"] > 0,
        "filename": filename,
        "source_type": source_type,
        "mode": mode,
        "doc_hash": doc_hash,
        "pages_loaded": len(pages),
        "chunks_total": insert_result["total"],
        "chunks_inserted": insert_result["inserted"],
        "chunks_deleted": deleted_count,
        "insert_errors": insert_result["errors"],
        "table_fqn": fqn_table,
        "index": index_full_name,
        "sync_triggered": sync_triggered,
        "sync_pending": sync_pending,
        "sync_error": sync_error,
        "duplicate_info": dup_info if dup_info["duplicate_type"] else None,
    }
